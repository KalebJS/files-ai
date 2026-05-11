"""Integration tests for mixed file-type extraction."""

from __future__ import annotations

import io
import os
import zipfile
from pathlib import Path

import pytest

from files_ai.extract import extract_file
from files_ai.storage import FileRef
from files_ai.storage import LocalFiles

LIVE_VISION_ENV = "RUN_LIVE_OLLAMA_EXTRACTION_TESTS"


def test_extract_common_file_types_integration(tmp_path: Path) -> None:
    """Extract text from common supported file types end-to-end."""
    files = LocalFiles(tmp_path)
    drop = FileRef("local", "/dropzone")
    files.make_dir(drop)

    csv_ref = files.join(drop, "budget.csv")
    (tmp_path / "dropzone" / "budget.csv").write_text(
        "Category,Amount\nFood,20\nRent,900\n",
        encoding="utf-8",
    )

    xlsx_ref = files.join(drop, "workbook.xlsx")
    (tmp_path / "dropzone" / "workbook.xlsx").write_bytes(_build_minimal_xlsx())

    docx_ref = files.join(drop, "notes.docx")
    _write_docx(tmp_path / "dropzone" / "notes.docx", "Quarterly planning memo")

    pptx_ref = files.join(drop, "slides.pptx")
    _write_pptx(tmp_path / "dropzone" / "slides.pptx", "Project Sunrise kickoff")

    pdf_ref = files.join(drop, "letter.pdf")
    _write_text_pdf(tmp_path / "dropzone" / "letter.pdf", "Tax planning checklist")

    pages_ref = files.join(drop, "notes.pages")
    _write_pages_with_index_xml(
        tmp_path / "dropzone" / "notes.pages",
        "Neighborhood budget draft",
    )

    eml_ref = files.join(drop, "message.eml")
    (tmp_path / "dropzone" / "message.eml").write_text(
        "Subject: Summer Plan\nFrom: user@example.com\nTo: other@example.com\n\n"
        "Discuss quarterly budget.",
        encoding="utf-8",
    )

    csv_result = extract_file(files, csv_ref, max_bytes=4096, vision_enabled=False)
    xlsx_result = extract_file(files, xlsx_ref, max_bytes=4096, vision_enabled=False)
    docx_result = extract_file(files, docx_ref, max_bytes=4096, vision_enabled=False)
    pptx_result = extract_file(files, pptx_ref, max_bytes=4096, vision_enabled=False)
    pdf_result = extract_file(files, pdf_ref, max_bytes=4096, vision_enabled=False)
    pages_result = extract_file(files, pages_ref, max_bytes=4096, vision_enabled=False)
    eml_result = extract_file(files, eml_ref, max_bytes=4096, vision_enabled=False)

    assert "Category | Amount" in csv_result.text
    assert "Sheet: Sheet1" in xlsx_result.text
    assert "Kaleb | 100" in xlsx_result.text
    assert "Quarterly planning memo" in docx_result.text
    assert "Project Sunrise kickoff" in pptx_result.text
    assert "Tax planning checklist" in pdf_result.text
    assert "Neighborhood budget draft" in pages_result.text
    assert "Subject: Summer Plan" in eml_result.text


def test_extract_image_and_scanned_pdf_with_live_vision(tmp_path: Path) -> None:
    """Optionally run live vision extraction against Ollama."""
    if os.getenv(LIVE_VISION_ENV, "").strip().lower() not in {"1", "true", "yes"}:
        pytest.skip(f"Set {LIVE_VISION_ENV}=true to run live vision extraction tests.")

    files = LocalFiles(tmp_path)
    drop = FileRef("local", "/dropzone")
    files.make_dir(drop)

    image_path = tmp_path / "dropzone" / "receipt.png"
    _write_png_with_text(image_path, "Invoice 1042 paid")
    image_ref = files.join(drop, "receipt.png")

    scanned_pdf_path = tmp_path / "dropzone" / "scan.pdf"
    _write_scanned_pdf(scanned_pdf_path, image_path)
    scanned_pdf_ref = files.join(drop, "scan.pdf")

    base_url = os.getenv("OLLAMA_BASE_URL", "http://host.docker.internal:11434")
    model = os.getenv("VISION_MODEL", os.getenv("MODEL", "gemma4:31b-cloud"))
    api_key = os.getenv("OLLAMA_API_KEY", "")

    image_result = extract_file(
        files,
        image_ref,
        max_bytes=4096,
        vision_enabled=True,
        vision_model=model,
        vision_base_url=base_url,
        vision_api_key=api_key,
    )
    scanned_pdf_result = extract_file(
        files,
        scanned_pdf_ref,
        max_bytes=4096,
        vision_enabled=True,
        vision_model=model,
        vision_base_url=base_url,
        vision_api_key=api_key,
    )

    assert image_result.tier == 3
    assert scanned_pdf_result.tier == 3
    assert image_result.text.strip()
    assert scanned_pdf_result.text.strip()
    assert image_result.text.strip() != "receipt.png"
    assert scanned_pdf_result.text.strip() != "scan.pdf"


def _write_docx(path: Path, text: str) -> None:
    """Write a minimal DOCX document."""
    from docx import Document

    document = Document()
    document.add_paragraph(text)
    document.save(path)


def _write_pptx(path: Path, text: str) -> None:
    """Write a minimal PPTX file with a title slide."""
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
    presentation.save(path)


def _write_text_pdf(path: Path, text: str) -> None:
    """Write a minimal text PDF file."""
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    document.save(path)
    document.close()


def _write_pages_with_index_xml(path: Path, text: str) -> None:
    """Write a minimal iWork-like archive with index.xml text."""
    payload = io.BytesIO()
    with zipfile.ZipFile(payload, "w") as archive:
        archive.writestr("index.xml", f"<root><p>{text}</p></root>")
    path.write_bytes(payload.getvalue())


def _write_png_with_text(path: Path, text: str) -> None:
    """Write a PNG image containing centered text."""
    from PIL import Image
    from PIL import ImageDraw

    image = Image.new("RGB", (900, 300), color="white")
    draw = ImageDraw.Draw(image)
    draw.text((40, 120), text, fill="black")
    image.save(path, format="PNG")


def _write_scanned_pdf(path: Path, image_path: Path) -> None:
    """Write an image-only PDF to exercise scanned-PDF vision flow."""
    import fitz

    image_bytes = image_path.read_bytes()
    document = fitz.open()
    page = document.new_page(width=900, height=300)
    page.insert_image(fitz.Rect(0, 0, 900, 300), stream=image_bytes)
    document.save(path)
    document.close()


def _build_minimal_xlsx() -> bytes:
    """Build a tiny XLSX payload for extraction integration tests."""
    payload = io.BytesIO()
    with zipfile.ZipFile(payload, "w") as archive:
        archive.writestr(
            "xl/workbook.xml",
            (
                '<?xml version="1.0"?>'
                "<workbook xmlns:r="
                '"http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
                ' xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
                '<sheets><sheet name="Sheet1" r:id="rId1"/></sheets>'
                "</workbook>"
            ),
        )
        archive.writestr(
            "xl/_rels/workbook.xml.rels",
            (
                '<?xml version="1.0"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="worksheet" '
                'Target="worksheets/sheet1.xml"/>'
                "</Relationships>"
            ),
        )
        archive.writestr(
            "xl/sharedStrings.xml",
            (
                '<?xml version="1.0"?>'
                '<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
                "<si><t>Name</t></si><si><t>Score</t></si><si><t>Kaleb</t></si>"
                "</sst>"
            ),
        )
        archive.writestr(
            "xl/worksheets/sheet1.xml",
            (
                '<?xml version="1.0"?>'
                '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
                "<sheetData>"
                '<row r="1"><c r="A1" t="s"><v>0</v></c>'
                '<c r="B1" t="s"><v>1</v></c></row>'
                '<row r="2"><c r="A2" t="s"><v>2</v></c>'
                '<c r="B2"><v>100</v></c></row>'
                "</sheetData>"
                "</worksheet>"
            ),
        )
    return payload.getvalue()
