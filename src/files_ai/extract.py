"""Tiered file-content extraction helpers."""

from __future__ import annotations

import base64
import csv
import io
import json
import mimetypes
import re
import tarfile
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from dataclasses import dataclass
from email import policy
from email.parser import BytesParser

from .storage import FileRef
from .storage import Files

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".tsv",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".log",
    ".ini",
    ".cfg",
    ".toml",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".java",
    ".go",
    ".rs",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".sh",
    ".sql",
}
IMAGE_EXTENSIONS = {
    ".jpg",
    ".jpeg",
    ".png",
    ".webp",
    ".bmp",
    ".gif",
    ".tif",
    ".tiff",
    ".heic",
    ".heif",
}
SPREADSHEET_EXTENSIONS = {
    ".xlsx",
    ".xlsm",
    ".xlsb",
    ".xls",
    ".ods",
    ".csv",
    ".tsv",
}
PANDOC_EXTENSIONS = {
    ".docx",
    ".odt",
    ".rtf",
    ".epub",
    ".html",
    ".htm",
    ".doc",
}
IWORK_EXTENSIONS = {".pages", ".key", ".keynote", ".numbers"}
ARCHIVE_EXTENSIONS = {".zip", ".tar", ".tgz", ".tar.gz"}
MIN_TEXT_THRESHOLD = 48
MAX_VISION_IMAGES = 3
VISION_IMAGE_MAX_SIDE = 1600
GENERIC_NAME = re.compile(r"^(img|scan|document|file)[-_ ]?\d*$", re.IGNORECASE)


@dataclass(frozen=True)
class ExtractionResult:
    """Normalized extraction output for one file.

    Attributes:
        text: Extracted text or fallback filename.
        mime: Detected MIME type when available.
        tier: Extraction tier used to produce `text`.
        encrypted: Whether the file is an encrypted PDF requiring a password.
    """

    text: str
    mime: str | None
    tier: int
    encrypted: bool = False


@dataclass(frozen=True)
class VisionConfig:
    """Vision extraction runtime configuration."""

    enabled: bool
    model: str
    base_url: str
    api_key: str


def extract_file(
    files: Files,
    ref: FileRef,
    *,
    max_bytes: int = 8192,
    vision_enabled: bool = True,
    vision_model: str = "gemma4:31b-cloud",
    vision_base_url: str = "https://ollama.com",
    vision_api_key: str = "",
) -> ExtractionResult:
    """Extract useful text from a file using tiered strategies.

    Args:
        files: Storage backend used for file reads.
        ref: File to extract text from.
        max_bytes: Maximum bytes to read for text-like content.
        vision_enabled: Whether vision extraction can be used when needed.
        vision_model: Vision-capable model name.
        vision_base_url: Ollama endpoint URL.
        vision_api_key: Optional Ollama API key.

    Returns:
        ExtractionResult: Normalized extraction result.
    """
    name = files.name_of(ref)
    first_bytes = files.read_bytes(ref, limit=min(max_bytes, 4096))
    mime = _detect_mime(name, first_bytes)
    suffix = _normalized_suffix(name.lower())
    tier = 1
    text = ""
    used_vision = False
    encrypted = False
    vision = VisionConfig(
        enabled=vision_enabled,
        model=vision_model,
        base_url=vision_base_url,
        api_key=vision_api_key,
    )

    if _needs_tier_two(name, mime):
        text, used_vision = _extract_text(
            files,
            ref,
            mime,
            suffix=suffix,
            max_bytes=max_bytes,
            vision=vision,
        )
        text = text.strip()
        if text:
            tier = 3 if used_vision else 2
        elif suffix == ".pdf" or mime == "application/pdf":
            encrypted = _is_pdf_encrypted(files.read_bytes(ref))

    if not text:
        text = name
    return ExtractionResult(text=text, mime=mime, tier=tier, encrypted=encrypted)


def _detect_mime(filename: str, content: bytes) -> str | None:
    """Detect MIME type from bytes with filename fallback.

    Args:
        filename: Source filename used for fallback guessing.
        content: Initial bytes used for content-based detection.

    Returns:
        str | None: MIME type when detected, otherwise `None`.
    """
    try:
        import magic  # type: ignore

        detected = magic.from_buffer(content, mime=True)
        if detected:
            return str(detected)
    except Exception:
        pass
    guessed, _ = mimetypes.guess_type(filename)
    return guessed


def _needs_tier_two(name: str, mime: str | None) -> bool:
    """Return whether richer text extraction is needed.

    Args:
        name: Source filename.
        mime: Detected MIME type.

    Returns:
        bool: `True` when tier-two extraction should run.
    """
    stem = name.rsplit(".", 1)[0]
    suffix = _normalized_suffix(name.lower())
    if GENERIC_NAME.match(stem):
        return True
    if suffix in TEXT_EXTENSIONS:
        return True
    if suffix in IMAGE_EXTENSIONS:
        return True
    if suffix in SPREADSHEET_EXTENSIONS:
        return True
    if suffix in PANDOC_EXTENSIONS:
        return True
    if suffix in IWORK_EXTENSIONS:
        return True
    if suffix in ARCHIVE_EXTENSIONS:
        return True
    if suffix in {".pdf", ".pptx", ".eml"}:
        return True
    if mime is None:
        return True
    if mime.startswith("text/"):
        return True
    return mime in {
        "application/pdf",
        "message/rfc822",
        "application/json",
        "application/xml",
        "application/zip",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    } or mime.startswith("image/")


def _extract_text(
    files: Files,
    ref: FileRef,
    mime: str | None,
    *,
    suffix: str,
    max_bytes: int,
    vision: VisionConfig,
) -> tuple[str, bool]:
    """Extract text from supported document formats.

    Args:
        files: Storage backend used for file reads.
        ref: File reference to extract from.
        mime: Detected MIME type.
        suffix: Lowercase filename suffix.
        max_bytes: Maximum bytes to read for text media.
        vision: Vision extraction settings.

    Returns:
        tuple[str, bool]: Extracted text and whether vision extraction was used.
    """
    if suffix in {".csv", ".tsv"}:
        return _extract_text_like(files, ref, max_bytes=max_bytes), False
    if _is_text_like(suffix=suffix, mime=mime):
        return _extract_text_like(files, ref, max_bytes=max_bytes), False
    payload = files.read_bytes(ref)
    if suffix == ".pdf" or mime == "application/pdf":
        return _extract_pdf_bytes(payload, max_bytes=max_bytes, vision=vision)
    if suffix in IMAGE_EXTENSIONS or (mime and mime.startswith("image/")):
        return _extract_image_bytes(payload, mime=mime, suffix=suffix, vision=vision)
    if suffix in SPREADSHEET_EXTENSIONS:
        sheet_text = _extract_spreadsheet_bytes(payload, suffix=suffix)
        if sheet_text:
            return sheet_text, False
    if suffix == ".pptx":
        ppt_text = _extract_pptx_bytes(payload)
        if ppt_text:
            return ppt_text, False
    if suffix in PANDOC_EXTENSIONS:
        doc_text = _extract_document_with_pandoc(payload, suffix=suffix)
        if doc_text:
            return doc_text, False
        if suffix == ".docx":
            fallback_docx = _extract_docx_with_python_docx(payload)
            if fallback_docx:
                return fallback_docx, False
    if suffix in IWORK_EXTENSIONS:
        return _extract_iwork_bytes(payload, max_bytes=max_bytes, vision=vision)
    if suffix == ".eml" or mime == "message/rfc822":
        eml_text = _extract_eml_bytes(payload, max_bytes=max_bytes)
        if eml_text:
            return eml_text, False
    if _is_archive(suffix):
        archive_text = _extract_archive_listing(payload, suffix=suffix)
        if archive_text:
            return archive_text, False
    if mime and mime.startswith("text/"):
        return files.read_bytes(ref, limit=max_bytes).decode(
            "utf-8", errors="ignore"
        ), (False)
    return "", False


def _is_text_like(*, suffix: str, mime: str | None) -> bool:
    """Return whether the file should be decoded as text directly."""
    if suffix in IMAGE_EXTENSIONS:
        return False
    if suffix in SPREADSHEET_EXTENSIONS:
        return False
    if suffix in PANDOC_EXTENSIONS:
        return False
    if suffix in IWORK_EXTENSIONS:
        return False
    if suffix in ARCHIVE_EXTENSIONS:
        return False
    if suffix in {".pdf", ".pptx"}:
        return False
    if suffix in TEXT_EXTENSIONS:
        return True
    if mime and mime.startswith("text/"):
        return True
    return mime in {"application/json", "application/xml"}


def _normalized_suffix(name: str) -> str:
    """Return normalized filename suffix with `.tar.gz` handling."""
    if name.endswith(".tar.gz"):
        return ".tar.gz"
    last = name.rsplit(".", 1)
    if len(last) != 2:
        return ""
    return f".{last[1]}"


def _extract_text_like(files: Files, ref: FileRef, *, max_bytes: int) -> str:
    """Read UTF-8 text content with bounded size."""
    payload = files.read_bytes(ref, limit=max_bytes)
    text = payload.decode("utf-8", errors="ignore")
    lower = files.name_of(ref).lower()
    if lower.endswith(".csv") or lower.endswith(".tsv"):
        delimiter = "," if lower.endswith(".csv") else "\t"
        lines: list[str] = []
        for index, row in enumerate(csv.reader(io.StringIO(text), delimiter=delimiter)):
            if index >= 120:
                break
            compact = [cell.strip() for cell in row if cell.strip()]
            if compact:
                lines.append(" | ".join(compact[:20]))
        return "\n".join(lines)
    if lower.endswith(".json"):
        try:
            return json.dumps(json.loads(text), indent=2)[: max_bytes * 2]
        except Exception:
            return text
    return text


def _extract_pdf_bytes(
    payload: bytes,
    *,
    max_bytes: int,
    vision: VisionConfig,
) -> tuple[str, bool]:
    """Extract PDF text, falling back to vision for scanned pages."""
    try:
        import fitz
    except Exception:
        return "", False
    try:
        with fitz.open(stream=payload, filetype="pdf") as document:
            if document.needs_pass or document.is_encrypted:
                return "", False
            page_text: list[str] = []
            for page in document:
                text = page.get_text("text").strip()
                if text:
                    page_text.append(text)
            joined = "\n\n".join(page_text).strip()
            if len(joined) >= MIN_TEXT_THRESHOLD:
                return joined[: max_bytes * 6], False
            if not vision.enabled:
                return joined, False
            images: list[tuple[bytes, str]] = []
            for page in list(document)[:MAX_VISION_IMAGES]:
                pixmap = page.get_pixmap(dpi=170)
                images.append((pixmap.tobytes("png"), "image/png"))
            vision_text = _extract_with_vision(
                images,
                vision=vision,
                prompt=(
                    "Extract readable document text. Keep line breaks and key fields."
                ),
            )
            if vision_text:
                return vision_text, True
            return joined, False
    except Exception:
        return "", False


def _is_pdf_encrypted(payload: bytes) -> bool:
    """Return whether a PDF payload requires a password to open pages."""
    try:
        import fitz
    except Exception:
        return False
    try:
        with fitz.open(stream=payload, filetype="pdf") as document:
            return bool(document.needs_pass or document.is_encrypted)
    except Exception:
        return False


def _extract_image_bytes(
    payload: bytes,
    *,
    mime: str | None,
    suffix: str,
    vision: VisionConfig,
) -> tuple[str, bool]:
    """Extract text from an image via vision model."""
    if not vision.enabled:
        return "", False
    image_mime = mime or _guess_image_mime(suffix)
    if image_mime is None:
        return "", False
    text = _extract_with_vision(
        [(payload, image_mime)],
        vision=vision,
        prompt=(
            "Extract all visible text and summarize document/image context in "
            "plain text."
        ),
    )
    if not text:
        return "", False
    return text, True


def _extract_with_vision(
    images: list[tuple[bytes, str]],
    *,
    vision: VisionConfig,
    prompt: str,
) -> str:
    """Send image payloads to a vision model and return extracted text."""
    if not images:
        return ""
    content: list[dict[str, str | dict[str, str]]] = [{"type": "text", "text": prompt}]
    for payload, mime in images[:MAX_VISION_IMAGES]:
        prepared = _prepare_vision_image(payload, mime=mime)
        if prepared is None:
            continue
        image_bytes, image_mime = prepared
        encoded = base64.b64encode(image_bytes).decode("utf-8")
        content.append(
            {
                "type": "image_url",
                "image_url": {"url": f"data:{image_mime};base64,{encoded}"},
            }
        )
    if len(content) == 1:
        return ""
    headers: dict[str, str] = {}
    if vision.api_key:
        headers["Authorization"] = f"Bearer {vision.api_key}"
    try:
        from langchain_core.messages import HumanMessage
        from langchain_ollama import ChatOllama
    except Exception:
        return ""
    try:
        llm = ChatOllama(
            model=vision.model,
            base_url=vision.base_url,
            client_kwargs={"headers": headers} if headers else None,
            temperature=0,
        )
        response = llm.invoke([HumanMessage(content=content)])
    except Exception:
        return ""
    return _coerce_llm_text(response).strip()


def _coerce_llm_text(response: object) -> str:
    """Normalize LLM response content into plain text."""
    content = getattr(response, "content", "")
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        lines: list[str] = []
        for block in content:
            if isinstance(block, str):
                lines.append(block)
            elif isinstance(block, dict):
                text = block.get("text")
                if isinstance(text, str):
                    lines.append(text)
        return "\n".join(line for line in lines if line.strip())
    return str(response)


def _prepare_vision_image(payload: bytes, *, mime: str) -> tuple[bytes, str] | None:
    """Resize large images before passing to the vision model."""
    if len(payload) <= 1_500_000:
        return payload, mime
    try:
        from PIL import Image
    except Exception:
        return None
    try:
        with Image.open(io.BytesIO(payload)) as image:
            image = image.convert("RGB")
            image.thumbnail((VISION_IMAGE_MAX_SIDE, VISION_IMAGE_MAX_SIDE))
            out = io.BytesIO()
            image.save(out, format="JPEG", quality=82, optimize=True)
            return out.getvalue(), "image/jpeg"
    except Exception:
        return None


def _guess_image_mime(suffix: str) -> str | None:
    """Guess image MIME type from filename suffix."""
    guessed, _ = mimetypes.guess_type(f"file{suffix}")
    return guessed


def _extract_document_with_pandoc(payload: bytes, *, suffix: str) -> str:
    """Extract document text via pypandoc when available."""
    try:
        import pypandoc
    except Exception:
        return ""
    try:
        with tempfile.NamedTemporaryFile(suffix=suffix) as temp:
            temp.write(payload)
            temp.flush()
            text = pypandoc.convert_file(temp.name, "plain")
            return text.strip()
    except Exception:
        return ""


def _extract_docx_with_python_docx(payload: bytes) -> str:
    """Fallback docx parser for environments without Pandoc."""
    try:
        from docx import Document
    except Exception:
        return ""
    try:
        document = Document(io.BytesIO(payload))
        return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()
    except Exception:
        return ""


def _extract_pptx_bytes(payload: bytes) -> str:
    """Extract text from PowerPoint `.pptx` payload."""
    try:
        from pptx import Presentation
    except Exception:
        return ""
    try:
        presentation = Presentation(io.BytesIO(payload))
    except Exception:
        return ""
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                lines.append(text.strip())
    return "\n".join(lines).strip()


def _extract_spreadsheet_bytes(payload: bytes, *, suffix: str) -> str:
    """Extract spreadsheet text via calamine or XML fallback."""
    text = _extract_spreadsheet_with_calamine(payload)
    if text:
        return text
    if suffix in {".xlsx", ".xlsm"}:
        return _extract_xlsx_xml(payload)
    return ""


def _extract_spreadsheet_with_calamine(payload: bytes) -> str:
    """Best-effort spreadsheet extraction using python-calamine."""
    try:
        import python_calamine as calamine
    except Exception:
        return ""
    workbook = _open_calamine_workbook(calamine, payload)
    if workbook is None:
        return ""
    sheet_names = getattr(workbook, "sheet_names", None)
    if callable(sheet_names):
        sheet_names = sheet_names()
    if not isinstance(sheet_names, list):
        return ""
    lines: list[str] = []
    for sheet_name in sheet_names[:5]:
        rows = _read_calamine_rows(workbook, sheet_name)
        if not rows:
            continue
        lines.append(f"Sheet: {sheet_name}")
        for row in rows[:120]:
            compact = [str(cell).strip() for cell in row if str(cell).strip()]
            if compact:
                lines.append(" | ".join(compact[:20]))
    return "\n".join(lines).strip()


def _open_calamine_workbook(calamine: object, payload: bytes) -> object | None:
    """Try known workbook constructors from python-calamine."""
    workbook_cls = getattr(calamine, "CalamineWorkbook", None)
    if workbook_cls is not None:
        for method_name in ("from_filelike", "from_bytes"):
            method = getattr(workbook_cls, method_name, None)
            if callable(method):
                try:
                    return method(io.BytesIO(payload))
                except Exception:
                    continue
    for func_name in ("load_workbook", "open_workbook"):
        func = getattr(calamine, func_name, None)
        if callable(func):
            try:
                return func(io.BytesIO(payload))
            except Exception:
                continue
    return None


def _read_calamine_rows(workbook: object, sheet_name: str) -> list[list[object]]:
    """Read rows from one calamine sheet using known APIs."""
    for method_name in ("get_sheet_by_name", "sheet_by_name", "worksheet_range"):
        method = getattr(workbook, method_name, None)
        if not callable(method):
            continue
        try:
            sheet = method(sheet_name)
        except Exception:
            continue
        rows = _coerce_sheet_rows(sheet)
        if rows:
            return rows
    return []


def _coerce_sheet_rows(sheet: object) -> list[list[object]]:
    """Normalize a calamine sheet object into row-major values."""
    if isinstance(sheet, list):
        return [row if isinstance(row, list) else [row] for row in sheet]
    converter = getattr(sheet, "to_python", None)
    if callable(converter):
        try:
            rows = converter()
            if isinstance(rows, list):
                return [row if isinstance(row, list) else [row] for row in rows]
        except Exception:
            pass
    rows_attr = getattr(sheet, "rows", None)
    if isinstance(rows_attr, list):
        return [row if isinstance(row, list) else [row] for row in rows_attr]
    iterator = getattr(sheet, "iter_rows", None)
    if callable(iterator):
        try:
            return [list(row) for row in iterator()]
        except Exception:
            return []
    return []


def _extract_xlsx_xml(payload: bytes) -> str:
    """Extract text from XLSX XML structure without extra dependencies."""
    try:
        archive = zipfile.ZipFile(io.BytesIO(payload))
    except zipfile.BadZipFile:
        return ""
    with archive:
        shared_strings = _xlsx_shared_strings(archive)
        sheets = _xlsx_sheets(archive)
        if not sheets:
            return ""
        lines: list[str] = []
        for sheet_name, sheet_path in sheets[:5]:
            try:
                raw_sheet = archive.read(sheet_path)
            except KeyError:
                continue
            rows = _xlsx_sheet_rows(raw_sheet, shared_strings)
            if not rows:
                continue
            lines.append(f"Sheet: {sheet_name}")
            for row in rows[:120]:
                compact = [cell.strip() for cell in row if cell.strip()]
                if compact:
                    lines.append(" | ".join(compact[:20]))
        return "\n".join(lines).strip()


def _xlsx_shared_strings(archive: zipfile.ZipFile) -> list[str]:
    """Return shared strings from XLSX archive."""
    try:
        raw = archive.read("xl/sharedStrings.xml")
    except KeyError:
        return []
    try:
        root = ET.fromstring(raw)
    except ET.ParseError:
        return []
    values: list[str] = []
    for item in root.findall(".//{*}si"):
        parts = [node.text or "" for node in item.findall(".//{*}t")]
        values.append("".join(parts))
    return values


def _xlsx_sheets(archive: zipfile.ZipFile) -> list[tuple[str, str]]:
    """Return workbook sheet names and XML paths."""
    try:
        workbook_raw = archive.read("xl/workbook.xml")
        rels_raw = archive.read("xl/_rels/workbook.xml.rels")
    except KeyError:
        return []
    try:
        workbook = ET.fromstring(workbook_raw)
        rels = ET.fromstring(rels_raw)
    except ET.ParseError:
        return []
    rel_map: dict[str, str] = {}
    for rel in rels.findall(".//{*}Relationship"):
        rel_id = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        if rel_id and target:
            rel_map[rel_id] = target
    sheets: list[tuple[str, str]] = []
    rel_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for sheet in workbook.findall(".//{*}sheet"):
        name = sheet.attrib.get("name", "Sheet")
        rel_id = sheet.attrib.get(f"{{{rel_ns}}}id")
        if not rel_id:
            continue
        target = rel_map.get(rel_id)
        if not target:
            continue
        normalized = target.lstrip("/")
        if not normalized.startswith("xl/"):
            normalized = f"xl/{normalized}"
        sheets.append((name, normalized))
    return sheets


def _xlsx_sheet_rows(raw_sheet: bytes, shared_strings: list[str]) -> list[list[str]]:
    """Extract row values from XLSX worksheet XML."""
    try:
        root = ET.fromstring(raw_sheet)
    except ET.ParseError:
        return []
    rows: list[list[str]] = []
    for row in root.findall(".//{*}sheetData/{*}row"):
        values: list[str] = []
        for cell in row.findall("{*}c"):
            value = _xlsx_cell_value(cell, shared_strings)
            if value:
                values.append(value)
        if values:
            rows.append(values)
    return rows


def _xlsx_cell_value(cell: ET.Element, shared_strings: list[str]) -> str:
    """Return one XLSX cell value from XML."""
    cell_type = cell.attrib.get("t")
    inline = cell.find("{*}is/{*}t")
    if inline is not None and inline.text:
        return inline.text
    value = cell.find("{*}v")
    if value is None or value.text is None:
        return ""
    if cell_type == "s":
        try:
            return shared_strings[int(value.text)]
        except Exception:
            return ""
    return value.text


def _extract_iwork_bytes(
    payload: bytes,
    *,
    max_bytes: int,
    vision: VisionConfig,
) -> tuple[str, bool]:
    """Extract text from Apple iWork archives using previews and metadata."""
    try:
        archive = zipfile.ZipFile(io.BytesIO(payload))
    except zipfile.BadZipFile:
        return "", False
    with archive:
        preview_pdf = _read_first_existing(
            archive,
            [
                "QuickLook/Preview.pdf",
                "Preview.pdf",
            ],
        )
        if preview_pdf:
            return _extract_pdf_bytes(preview_pdf, max_bytes=max_bytes, vision=vision)
        preview_images = _read_existing_images(archive)
        if preview_images:
            text = _extract_with_vision(
                preview_images,
                vision=vision,
                prompt=(
                    "Extract readable text from this Apple iWork preview and "
                    "summarize key content."
                ),
            )
            if text:
                return text, True
        xml_payload = _read_first_existing(
            archive,
            ["index.xml", "index.apxl", "Index/Document.iwa"],
        )
        if xml_payload:
            try:
                xml_text = xml_payload.decode("utf-8", errors="ignore")
            except Exception:
                xml_text = ""
            if xml_text:
                plain = re.sub(r"<[^>]+>", " ", xml_text)
                plain = re.sub(r"\s+", " ", plain).strip()
                return plain[: max_bytes * 4], False
    return "", False


def _read_first_existing(
    archive: zipfile.ZipFile,
    names: list[str],
) -> bytes | None:
    """Return bytes for the first path that exists in the archive."""
    for name in names:
        try:
            return archive.read(name)
        except KeyError:
            continue
    return None


def _read_existing_images(archive: zipfile.ZipFile) -> list[tuple[bytes, str]]:
    """Read preview images from archive."""
    images: list[tuple[bytes, str]] = []
    for name in archive.namelist():
        lower = name.lower()
        if not lower.endswith(tuple(IMAGE_EXTENSIONS)):
            continue
        try:
            payload = archive.read(name)
        except KeyError:
            continue
        suffix = _normalized_suffix(lower)
        mime = _guess_image_mime(suffix)
        if mime:
            images.append((payload, mime))
    return images[:MAX_VISION_IMAGES]


def _extract_eml_bytes(payload: bytes, *, max_bytes: int) -> str:
    """Extract key fields and body text from `.eml` files."""
    try:
        message = BytesParser(policy=policy.default).parsebytes(payload)
    except Exception:
        return ""
    lines: list[str] = []
    for header in ("subject", "from", "to", "date"):
        value = message.get(header)
        if value:
            lines.append(f"{header.title()}: {value}")
    body = message.get_body(preferencelist=("plain",))
    if body is not None:
        try:
            text = body.get_content()
        except Exception:
            text = ""
        if text:
            lines.append("")
            lines.append(text.strip())
    return "\n".join(lines)[: max_bytes * 4].strip()


def _is_archive(suffix: str) -> bool:
    """Return whether suffix is a supported archive type."""
    return suffix in ARCHIVE_EXTENSIONS


def _extract_archive_listing(payload: bytes, *, suffix: str) -> str:
    """Return archive file listing summary."""
    if suffix == ".zip":
        try:
            with zipfile.ZipFile(io.BytesIO(payload)) as archive:
                names = archive.namelist()
        except zipfile.BadZipFile:
            return ""
    else:
        try:
            with tarfile.open(fileobj=io.BytesIO(payload), mode="r:*") as archive:
                names = [member.name for member in archive.getmembers()]
        except Exception:
            return ""
    if not names:
        return ""
    lines = ["Archive contents:"] + [f"- {name}" for name in names[:200]]
    return "\n".join(lines)
